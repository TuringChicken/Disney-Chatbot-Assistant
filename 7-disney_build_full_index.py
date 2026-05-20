# -*- coding: utf-8 -*-
"""
迪士尼RAG助手 - 全量知识库构建（按分类建立独立索引）

功能：
  扫描 Disney_RAG_KnowledgeBase/ 下 5 个子目录，
  每个子目录单独生成 FAISS 向量索引，
  同时合并生成全局统一索引 disney_full

支持格式:
  .docx / .doc  →  python-docx
  .pdf          →  PyMuPDF (pip install pymupdf)
  .pptx         →  python-pptx (pip install python-pptx)
  .ppt          →  跳过（需 LibreOffice 或 win32com）
  图片           →  multimodal embedding

运行: python 7-disney_build_full_index.py
"""
import os
import base64
import json
import time
import numpy as np
import faiss
import dashscope
from http import HTTPStatus
from pathlib import Path

# ── 可选依赖 ──────────────────────────────────────────────────────────────────
try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("[警告] python-docx 未安装，.docx/.doc 将被跳过。安装: pip install python-docx")

try:
    import fitz
    # 验证是否为真正的 PyMuPDF（而非同名占位包）
    _ = fitz.open  # noqa
    HAS_PDF = True
except (ImportError, AttributeError):
    HAS_PDF = False
    print("[警告] PyMuPDF 未安装，.pdf 将被跳过。安装: pip install pymupdf")

try:
    from pptx import Presentation as PptxPresentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("[警告] python-pptx 未安装，.pptx 将被跳过。安装: pip install python-pptx")

# ── 配置 ──────────────────────────────────────────────────────────────────────
DASHSCOPE_API_KEY = os.getenv("DASHSCOPE_API_KEY")
if not DASHSCOPE_API_KEY:
    raise ValueError("请设置环境变量 DASHSCOPE_API_KEY")
dashscope.api_key = DASHSCOPE_API_KEY

KB_ROOT = "Disney_RAG_KnowledgeBase"
OUTPUT_DIR = "disney_indexes"
MULTIMODAL_EMBEDDING_MODEL = "qwen3-vl-embedding"
CHUNK_SIZE = 500
CHUNK_OVERLAP = 50
IMAGE_EXTS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp'}

# 子目录 → (索引文件前缀, 人类可读分类名)
CATEGORIES = {
    "1-产品与服务详情":        ("cat1_products",    "产品与服务详情"),
    "2-运营流程与标准作业程序": ("cat2_operations",  "运营流程与标准作业程序"),
    "3-特殊情况与应急预案":     ("cat3_emergency",   "特殊情况与应急预案"),
    "4-客户关系与支持话术":     ("cat4_customer",    "客户关系与支持话术"),
    "5-内部知识与工具":        ("cat5_internal",    "内部知识与工具"),
}


# ── 文件解析 ───────────────────────────────────────────────────────────────────
def parse_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    parts = []
    for element in doc.element.body:
        if element.tag.endswith('p'):
            text = "".join(
                run.text for run in
                element.findall('.//w:t', {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'})
                if run.text
            ).strip()
            if text:
                parts.append(text)
        elif element.tag.endswith('tbl'):
            matched = [t for t in doc.tables if t._element is element]
            if not matched:
                continue
            table = matched[0]
            if not table.rows:
                continue
            header = [c.text.strip() for c in table.rows[0].cells]
            rows_md = ["| " + " | ".join(header) + " |", "|" + "---|" * len(header)]
            for row in table.rows[1:]:
                rows_md.append("| " + " | ".join(c.text.strip() for c in row.cells) + " |")
            parts.append("\n".join(rows_md))
    return "\n".join(parts)


def parse_pdf(file_path: str) -> str:
    doc = fitz.open(file_path)
    return "\n".join(page.get_text() for page in doc)


def parse_pptx(file_path: str) -> str:
    prs = PptxPresentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"[幻灯片 {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def split_text(text: str) -> list[str]:
    chunks, start = [], 0
    while start < len(text):
        chunk = text[start:start + CHUNK_SIZE].strip()
        if chunk:
            chunks.append(chunk)
        start += CHUNK_SIZE - CHUNK_OVERLAP
    return chunks


def extract_text(file_path: str) -> str | None:
    """根据扩展名分派到对应解析器，返回 None 表示跳过"""
    suffix = Path(file_path).suffix.lower()

    if suffix in ('.docx', '.doc'):
        if not HAS_DOCX:
            return None
        try:
            return parse_docx(file_path)
        except Exception as e:
            print(f"    [跳过] 解析失败 {Path(file_path).name}: {e}")
            return None

    if suffix == '.pdf':
        if not HAS_PDF:
            return None
        try:
            return parse_pdf(file_path)
        except Exception as e:
            print(f"    [跳过] 解析失败 {Path(file_path).name}: {e}")
            return None

    if suffix == '.pptx':
        if not HAS_PPTX:
            return None
        try:
            return parse_pptx(file_path)
        except Exception as e:
            print(f"    [跳过] 解析失败 {Path(file_path).name}: {e}")
            return None

    if suffix == '.ppt':
        print(f"    [跳过] .ppt 格式需要 LibreOffice，暂不支持: {Path(file_path).name}")
        return None

    return None


# ── Embedding（含重试） ────────────────────────────────────────────────────────
def _call_with_retry(fn, *args, max_retries=3, **kwargs):
    for attempt in range(max_retries):
        try:
            return fn(*args, **kwargs)
        except Exception as e:
            if attempt == max_retries - 1:
                raise
            wait = 2 ** attempt
            print(f"    [重试 {attempt+1}/{max_retries}] {e}，等待 {wait}s")
            time.sleep(wait)


def get_text_embedding(text: str) -> list:
    resp = _call_with_retry(
        dashscope.MultiModalEmbedding.call,
        model=MULTIMODAL_EMBEDDING_MODEL,
        input=[{"text": text}]
    )
    if resp.status_code != HTTPStatus.OK:
        raise Exception(f"文本 Embedding 失败: {resp.message}")
    return resp.output["embeddings"][0]["embedding"]


def get_image_embedding(image_path: str) -> list:
    with open(image_path, "rb") as f:
        b64 = base64.b64encode(f.read()).decode("utf-8")
    ext = Path(image_path).suffix.lower().lstrip(".")
    if ext == "jpg":
        ext = "jpeg"
    resp = _call_with_retry(
        dashscope.MultiModalEmbedding.call,
        model=MULTIMODAL_EMBEDDING_MODEL,
        input=[{"image": f"data:image/{ext};base64,{b64}"}]
    )
    if resp.status_code != HTTPStatus.OK:
        raise Exception(f"图片 Embedding 失败: {resp.message}")
    return resp.output["embeddings"][0]["embedding"]


# ── 工具 ──────────────────────────────────────────────────────────────────────
def should_skip(filename: str) -> bool:
    """跳过百度网盘未完成下载和隐藏文件"""
    return filename.startswith('.') or '.baiduyun' in filename or filename.endswith('.downloading')


def build_faiss(vectors: list, metadata: list, prefix: str, label: str):
    """将向量列表和元数据写入 FAISS 索引文件"""
    if not vectors:
        print(f"  [{label}] 无有效向量，跳过保存")
        return

    dim = len(vectors[0])
    index = faiss.IndexFlatL2(dim)
    index.add(np.array(vectors, dtype="float32"))

    idx_path = os.path.join(OUTPUT_DIR, f"{prefix}_index.faiss")
    meta_path = os.path.join(OUTPUT_DIR, f"{prefix}_metadata.json")

    faiss.write_index(index, idx_path)
    with open(meta_path, "w", encoding="utf-8") as f:
        json.dump(metadata, f, ensure_ascii=False, indent=2)

    text_n = sum(1 for m in metadata if m["type"] == "text")
    img_n = sum(1 for m in metadata if m["type"] == "image")
    print(f"  [{label}] 保存完成 → {prefix}_index.faiss  "
          f"(文本块:{text_n}, 图片:{img_n}, 总计:{len(metadata)})")


# ── 处理单个子目录 ─────────────────────────────────────────────────────────────
def process_directory(dirpath: str, category_name: str) -> tuple[list, list]:
    """
    扫描目录下所有文件，返回 (vectors, metadata_list)
    会递归处理子目录中的图片（如 images/ 子目录）
    """
    vectors, metadata = [], []
    doc_id = 0

    def _process_file(fpath: str):
        nonlocal doc_id
        fname = Path(fpath).name
        suffix = Path(fpath).suffix.lower()

        if should_skip(fname):
            return

        # ── 图片 ──────────────────────────────
        if suffix in IMAGE_EXTS:
            print(f"    [图片] {fname}")
            try:
                vec = get_image_embedding(fpath)
                metadata.append({
                    "id": doc_id,
                    "category": category_name,
                    "source": fname,
                    "type": "image",
                    "path": fpath,
                    "content": f"[图片] {fname}"
                })
                vectors.append(vec)
                doc_id += 1
            except Exception as e:
                print(f"    [跳过] 图片 embedding 失败 {fname}: {e}")
            return

        # ── 文档 ──────────────────────────────
        text = extract_text(fpath)
        if text is None:
            return

        chunks = split_text(text)
        if not chunks:
            print(f"    [跳过] {fname} 无有效文本")
            return

        print(f"    [文档] {fname}  →  {len(chunks)} 个 chunk")
        for chunk in chunks:
            try:
                vec = get_text_embedding(chunk)
                metadata.append({
                    "id": doc_id,
                    "category": category_name,
                    "source": fname,
                    "type": "text",
                    "content": chunk
                })
                vectors.append(vec)
                doc_id += 1
            except Exception as e:
                print(f"    [跳过] chunk embedding 失败: {e}")

    # 遍历目录（含一级子目录，用于处理 images/ 等）
    for entry in sorted(os.scandir(dirpath), key=lambda e: e.name):
        if entry.is_file():
            _process_file(entry.path)
        elif entry.is_dir():
            print(f"  → 扫描子目录: {entry.name}/")
            for sub_entry in sorted(os.scandir(entry.path), key=lambda e: e.name):
                if sub_entry.is_file():
                    _process_file(sub_entry.path)

    return vectors, metadata


# ── 主流程 ────────────────────────────────────────────────────────────────────
def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    all_vectors: list = []
    all_metadata: list = []
    global_id = 0

    print(f"\n{'='*60}")
    print(f"  迪士尼全量知识库构建  |  来源: {KB_ROOT}/")
    print(f"  输出目录: {OUTPUT_DIR}/")
    print(f"{'='*60}\n")

    # ── 处理根目录下的概览文档 ────────────────────────────────
    root_doc = os.path.join(KB_ROOT, "迪士尼RAG知识库.docx")
    if os.path.exists(root_doc):
        print("── 根目录概览文档 ─────────────────────────────────")
        text = extract_text(root_doc)
        if text:
            chunks = split_text(text)
            print(f"  迪士尼RAG知识库.docx  →  {len(chunks)} 个 chunk")
            for chunk in chunks:
                try:
                    vec = get_text_embedding(chunk)
                    meta = {
                        "id": global_id,
                        "category": "概览",
                        "source": "迪士尼RAG知识库.docx",
                        "type": "text",
                        "content": chunk
                    }
                    all_vectors.append(vec)
                    all_metadata.append(meta)
                    global_id += 1
                except Exception as e:
                    print(f"  [跳过] chunk 失败: {e}")
        print()

    # ── 按分类处理子目录 ──────────────────────────────────────
    for dirname, (prefix, label) in CATEGORIES.items():
        dirpath = os.path.join(KB_ROOT, dirname)
        if not os.path.isdir(dirpath):
            print(f"[警告] 目录不存在，跳过: {dirpath}")
            continue

        print(f"── {label} ({'─'*(45-len(label))})")
        vectors, metadata = process_directory(dirpath, label)

        # 保存当前分类的独立索引
        build_faiss(vectors, metadata, prefix, label)
        print()

        # 合并到全局索引（重新分配 id）
        for meta, vec in zip(metadata, vectors):
            m = dict(meta)
            m["id"] = global_id
            all_metadata.append(m)
            all_vectors.append(vec)
            global_id += 1

    # ── 保存全局统一索引 ──────────────────────────────────────
    print("── 全局统一索引 ───────────────────────────────────────")
    build_faiss(all_vectors, all_metadata, "disney_full", "全局")

    # ── 汇总统计 ─────────────────────────────────────────────
    print(f"\n{'='*60}")
    print("  构建完成！各分类统计：")
    print(f"{'='*60}")
    for dirname, (prefix, label) in CATEGORIES.items():
        cat_meta = [m for m in all_metadata if m.get("category") == label]
        if cat_meta:
            t = sum(1 for m in cat_meta if m["type"] == "text")
            i = sum(1 for m in cat_meta if m["type"] == "image")
            print(f"  {label:<22}  文本块:{t:4d}  图片:{i:3d}  合计:{t+i:4d}")
    print(f"  {'─'*50}")
    t_all = sum(1 for m in all_metadata if m["type"] == "text")
    i_all = sum(1 for m in all_metadata if m["type"] == "image")
    print(f"  {'全局合计':<22}  文本块:{t_all:4d}  图片:{i_all:3d}  合计:{t_all+i_all:4d}")
    print(f"{'='*60}\n")
    print(f"索引文件保存在: {os.path.abspath(OUTPUT_DIR)}/")


if __name__ == "__main__":
    main()
